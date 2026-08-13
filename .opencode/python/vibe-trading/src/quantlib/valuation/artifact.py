"""Versioned model artifacts for DCF / comps / three-statement, plus export.

:mod:`dcf`, :mod:`comps` and :mod:`threestatement` each answer "what is this
worth" for one run. None of them answer three questions every real valuation
practice needs answered on top of that: *which* inputs produced this number,
*can I tell if someone changed one*, and *can I hand this to someone who was
not in the room*. This module exists only for those three questions -- it does
not recompute anything the three model modules already compute, and it never
imports anything from them that isn't already public (:class:`DCFResult`,
:class:`CompsResult`, :class:`ThreeStatementProjection` and their nested
dataclasses, plus :class:`~.contracts.Assumption`).

THE ARTIFACT
------------
:func:`build_dcf_artifact`, :func:`build_comps_artifact` and
:func:`build_three_statement_artifact` each wrap one completed model run (the
raw inputs the caller fed the model, plus the call-level knobs such as
``discounting_convention`` that also determine the output, plus the model's
own result object) into one :class:`ModelArtifact`. Four things happen on the
way in:

1. **Input hash.** Every input -- the raw mapping/dataclasses the caller
   passed AND the call-level configuration knobs, because a knob like
   ``discounting_convention`` changes the output just as much as a number in
   ``inputs`` does -- is walked into a flat set of ``(dotted.path, value)``
   leaves (:func:`_flatten`), each leaf's number canonicalised
   (:func:`_normalize_number`: fixed 12-significant-digit scientific
   notation, ``-0.0`` collapsed to ``0.0``, ``int``/``float`` treated as the
   same economic value), then hashed with ``sha256`` over a
   ``sort_keys=True`` JSON encoding so the result depends on the *content* of
   the inputs and nothing about how they were typed in (dict key order,
   ``int`` vs ``float``, ``-0.0`` vs ``0.0``). Change one number anywhere in
   the input tree and the hash changes; two calls with the same inputs in a
   different key order produce the identical hash. See :func:`compute_input_hash`.
2. **Model version.** :func:`_module_version` hashes the producing module's
   own source file together with :mod:`.contracts`'s (every model in this
   package depends on ``contracts.py``'s input discipline, so a change there
   is a change to what produced the artifact even if the model file's own
   text is untouched). This is content-addressed on purpose: there is no
   ``__version__`` anywhere in :mod:`src.quantlib` to defer to (checked before
   writing this module), and a hand-maintained version string goes stale the
   first time someone edits the model and forgets to bump it. A content hash
   cannot go stale.
3. **Assumptions.** The same flatten pass that builds the input hash also
   walks the *result* object (not the raw inputs -- see
   :func:`build_dcf_artifact`'s docstring for why) collecting every
   :class:`~.contracts.Assumption` it finds, with the dotted path it was found
   at. Only :mod:`dcf` requires callers to pass :class:`Assumption` objects
   (``terminal_growth``, ``exit_multiple``); :mod:`comps` and
   :mod:`threestatement` take bare numbers for everything and therefore
   produce artifacts with an empty ``assumptions`` tuple. That is not a bug in
   this module -- it is an accurate report of what those two models do and do
   not ask the caller to justify, and every export below says so explicitly
   rather than leaving the section blank.
4. **Timestamp, supplied not invented.** ``generated_at`` is a required
   keyword argument of ``datetime`` type with no default anywhere in this
   file. This module never calls ``datetime.now()`` or ``time.time()`` --
   omitting the argument is a ``TypeError`` from Python's own argument
   binding, and passing a naive (non-timezone-aware) datetime is a
   ``ValueError`` from :func:`_require_generated_at`. A caller who wants a
   real-time artifact passes ``datetime.now(timezone.utc)`` themselves; this
   module has no opinion on "now" and cannot silently manufacture one, which
   is exactly what makes ``input_hash`` and every export below reproducible
   in a test.

DIFFING
-------
:func:`diff_artifacts` compares two artifacts of the same ``model_name`` by
set-comparing their flattened input and output leaf paths: paths present in
both but changed get a :class:`FieldChange` (with a numeric ``delta`` when
both sides are numbers), paths only on one side are reported separately as
additions/removals. This is the only way to answer "did the OUTPUT move
because an INPUT moved, or because the model code changed" -- which is why
``input_hash`` is computed from inputs and configuration only, never from
``model_version`` or ``generated_at``: two artifacts built from identical
inputs on different days, or by different code versions, must still hash
identically on the input side, or that comparison is not possible.

EXPORTS DO NOT HIDE THE UGLY PARTS
-----------------------------------
Every ``export_*`` function writes an "Assumptions & Data Quality" sheet (xlsx)
or a "Data Quality Notes" section (pptx) carrying, verbatim: every
:class:`~.contracts.Assumption` with its ``basis``, the ``input_hash`` and
``model_version`` that produced the run, and ``ModelArtifact.excluded`` --
which for :mod:`comps` includes every peer dropped from a multiple
(:class:`~.comps.ExcludedMultiple`) and every EV-bridge component a peer or
the target did not supply (:class:`~.comps.EVBridgeResult.omitted_components`),
and for :mod:`threestatement` includes every period where the revolver had to
draw or repay (the model's one declared plug) plus the convergence diagnostics
:class:`~.threestatement.PeriodResult` already carries. An export that only
shows the numbers that make the model look clean is not an export of the
model -- it is a different, dishonest document that happens to share some
numbers with it.

WHAT THIS MODULE DOES NOT DO
------------------------------
It does not validate valuation inputs (that is :mod:`.contracts`'s job,
already done by the time a result object exists) and it does not recompute a
sensitivity grid on the caller's behalf -- :func:`export_dcf_workbook`
requires the caller to pass one (e.g. from
:func:`src.quantlib.valuation.dcf.sensitivity_grid`) rather than inventing a
WACC/growth sweep range, because a sweep range is itself a judgement call
this package does not make silently for the same reason
:mod:`.contracts` refuses to default a terminal growth rate.
"""

from __future__ import annotations

import dataclasses
import hashlib
import json
import math
from collections.abc import Mapping, Sequence
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from types import MappingProxyType, ModuleType
from typing import Any

import numpy as np
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.text.text import TextFrame
from pptx.util import Inches, Pt

from src.quantlib.valuation import comps as _comps_module
from src.quantlib.valuation import contracts as _contracts_module
from src.quantlib.valuation import dcf as _dcf_module
from src.quantlib.valuation import threestatement as _threestatement_module
from src.quantlib.valuation.comps import (
    MULTIPLE_NAMES,
    CompsResult,
    PeerCompany,
    TargetCompany,
)
from src.quantlib.valuation.contracts import Assumption
from src.quantlib.valuation.dcf import DEFAULT_LONG_RUN_GDP_GROWTH_CEILING, DCFResult
from src.quantlib.valuation.threestatement import (
    BALANCE_TOLERANCE_REL,
    CIRCULARITY_TOLERANCE_REL,
    MAX_CIRCULARITY_ITERATIONS,
    ThreeStatementProjection,
)

__all__ = [
    "AssumptionRecord",
    "ModelArtifact",
    "FieldChange",
    "ArtifactDiff",
    "compute_input_hash",
    "build_dcf_artifact",
    "build_comps_artifact",
    "build_three_statement_artifact",
    "diff_artifacts",
    "export_dcf_workbook",
    "export_comps_workbook",
    "export_three_statement_workbook",
    "export_summary_deck",
]

#: Names this module knows how to build an artifact for. Any other
#: ``model_name`` on a :class:`ModelArtifact` is a construction bug, not a
#: valid but-unsupported case -- see :meth:`ModelArtifact.__post_init__`.
_MODEL_NAMES: tuple[str, ...] = ("dcf", "comps", "three_statement")

#: This module's own artifact-schema version. Bumped only if the normalization
#: scheme in :func:`_flatten` / :func:`_normalize_number` changes in a way
#: that would make an old and a new hash falsely comparable -- folded into
#: every input hash via :data:`_HASH_SCHEMA` so such a change is never silent.
_SCHEMA_VERSION = "artifact-v1"
_HASH_SCHEMA = "vt-valuation-input-hash-v1"


# ---------------------------------------------------------------------------
# Core types
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class AssumptionRecord:
    """One :class:`~.contracts.Assumption` found inside a result object, located.

    Attributes:
        path: Dotted/bracket path within the result object where this
            assumption was found, e.g. ``"$.terminal_value.terminal_growth"``.
        assumption: The :class:`~.contracts.Assumption` itself -- name, value,
            basis and optional source, unmodified.
    """

    path: str
    assumption: Assumption


@dataclass(frozen=True)
class ModelArtifact:
    """A reproducible, diffable snapshot of one completed valuation run.

    Attributes:
        model_name: One of :data:`_MODEL_NAMES` -- ``"dcf"``, ``"comps"`` or
            ``"three_statement"``.
        model_version: Content-addressed tag for the code that produced
            ``result`` -- see :func:`_module_version`.
        schema_version: This artifact schema's own version, :data:`_SCHEMA_VERSION`.
        generated_at: Caller-supplied, timezone-aware timestamp of the run.
            Never computed inside this module -- see the module docstring.
        input_hash: ``sha256`` of the canonicalised inputs and call-level
            configuration that produced ``result``. Does NOT depend on
            ``generated_at`` or ``model_version``, by construction -- see
            :func:`diff_artifacts`.
        inputs: Flat ``{dotted.path: value}`` snapshot of every input leaf
            (readable Python scalars, not the canonical hash strings), for
            display and diffing.
        outputs: Flat ``{dotted.path: value}`` snapshot of every leaf in
            ``result``, same convention as ``inputs``.
        assumptions: Every :class:`AssumptionRecord` found inside ``result``.
            Empty for ``comps`` and ``three_statement`` -- neither model's
            required inputs go through :class:`~.contracts.Assumption`; see
            the module docstring.
        result: The original typed result object (:class:`DCFResult`,
            :class:`CompsResult` or :class:`ThreeStatementProjection`),
            unmodified, for exporters that need typed access.
        excluded: Human-readable notes on everything the model dropped,
            omitted, or flagged rather than silently absorbed -- excluded
            comps, omitted EV-bridge components, terminal growth above the
            GDP-growth reference ceiling, revolver draws/repayments. Always
            present (an empty tuple IS the report "nothing was excluded"),
            and always written into every export -- see the module docstring.
        warnings: Warning strings the underlying result itself carries
            (currently only :attr:`CompsResult.warnings`); ``()`` for models
            that carry none.
    """

    model_name: str
    model_version: str
    schema_version: str
    generated_at: datetime
    input_hash: str
    inputs: Mapping[str, Any]
    outputs: Mapping[str, Any]
    assumptions: tuple[AssumptionRecord, ...]
    result: Any
    excluded: tuple[str, ...]
    warnings: tuple[str, ...]

    def __post_init__(self) -> None:
        """Check the artifact's own invariants.

        Raises:
            ValueError: If ``model_name`` is not one of :data:`_MODEL_NAMES`,
                or ``input_hash`` is not a 64-character hex sha256 digest.
        """
        if self.model_name not in _MODEL_NAMES:
            raise ValueError(
                f"ModelArtifact.model_name must be one of {_MODEL_NAMES}, "
                f"got {self.model_name!r}"
            )
        if len(self.input_hash) != 64 or any(
            c not in "0123456789abcdef" for c in self.input_hash
        ):
            raise ValueError(
                f"ModelArtifact.input_hash must be a 64-character sha256 hex "
                f"digest, got {self.input_hash!r}"
            )


@dataclass(frozen=True)
class FieldChange:
    """One field that differs between two artifacts' inputs or outputs.

    Attributes:
        path: The dotted/bracket path that changed.
        old: The baseline artifact's value at ``path``.
        new: The updated artifact's value at ``path``.
        delta: ``new - old`` when both are numbers (and neither is a
            ``bool``); ``None`` otherwise (e.g. a string or boolean field, or
            a field that changed type).
    """

    path: str
    old: Any
    new: Any
    delta: float | None


@dataclass(frozen=True)
class ArtifactDiff:
    """The difference between two :class:`ModelArtifact` instances of the same model.

    Attributes:
        model_name: The shared ``model_name`` of both artifacts.
        baseline_generated_at: ``baseline.generated_at``.
        updated_generated_at: ``updated.generated_at``.
        baseline_input_hash: ``baseline.input_hash``.
        updated_input_hash: ``updated.input_hash``.
        input_hash_changed: ``baseline_input_hash != updated_input_hash``.
        baseline_model_version: ``baseline.model_version``.
        updated_model_version: ``updated.model_version``.
        model_version_changed: ``baseline_model_version != updated_model_version``.
        changed_inputs: One :class:`FieldChange` per input path present on
            both artifacts whose value differs.
        added_input_paths: Input paths present on ``updated`` only.
        removed_input_paths: Input paths present on ``baseline`` only.
        changed_outputs: One :class:`FieldChange` per output path present on
            both artifacts whose value differs.
        added_output_paths: Output paths present on ``updated`` only.
        removed_output_paths: Output paths present on ``baseline`` only.
    """

    model_name: str
    baseline_generated_at: datetime
    updated_generated_at: datetime
    baseline_input_hash: str
    updated_input_hash: str
    input_hash_changed: bool
    baseline_model_version: str
    updated_model_version: str
    model_version_changed: bool
    changed_inputs: tuple[FieldChange, ...]
    added_input_paths: tuple[str, ...]
    removed_input_paths: tuple[str, ...]
    changed_outputs: tuple[FieldChange, ...]
    added_output_paths: tuple[str, ...]
    removed_output_paths: tuple[str, ...]


# ---------------------------------------------------------------------------
# Timestamp discipline -- no hidden clock anywhere in this module
# ---------------------------------------------------------------------------


def _require_generated_at(value: datetime, model: str) -> datetime:
    """Check that a caller-supplied timestamp is usable, and only that.

    This function never reads a clock. Its only job is to reject a call that
    did not supply one properly -- see the module docstring.

    Args:
        value: The candidate ``generated_at``.
        model: Name of the calling builder, for the error message.

    Returns:
        ``value`` unchanged.

    Raises:
        TypeError: If ``value`` is not a ``datetime.datetime``.
        ValueError: If ``value`` is timezone-naive.
    """
    if not isinstance(value, datetime):
        raise TypeError(
            f"{model}: generated_at must be a datetime.datetime supplied by the "
            f"caller, got {type(value).__name__}. This module never calls "
            "datetime.now() -- pass datetime.now(timezone.utc) yourself."
        )
    if value.tzinfo is None or value.utcoffset() is None:
        raise ValueError(
            f"{model}: generated_at must be timezone-aware (e.g. "
            "datetime.now(timezone.utc)), got a naive datetime "
            f"{value!r} -- an artifact's timestamp must be unambiguous."
        )
    return value


# ---------------------------------------------------------------------------
# Canonical normalization: numbers, then the recursive flatten/hash engine
# ---------------------------------------------------------------------------


def _normalize_number(value: float) -> str:
    """Canonicalise one number for hashing.

    Fixed 12-significant-digit scientific notation, independent of magnitude,
    so ``int`` and ``float`` inputs representing the same economic value hash
    identically, and ``-0.0`` collapses to ``0.0``. 12 significant digits is
    far beyond the precision any input in this package carries, so a genuine
    input change always moves the hash while floating-point representation
    noise never does.

    Args:
        value: An ``int`` or ``float`` (already unwrapped from any numpy
            scalar by the caller).

    Returns:
        A canonical string form -- ``"NaN"`` / ``"Infinity"`` / ``"-Infinity"``
        for non-finite values, else fixed-precision scientific notation.
    """
    numeric = float(value)
    if math.isnan(numeric):
        return "NaN"
    if math.isinf(numeric):
        return "Infinity" if numeric > 0 else "-Infinity"
    if numeric == 0.0:
        numeric = 0.0  # collapse -0.0
    return f"{numeric:.12e}"


def _flatten(
    value: Any,
    path: str,
    hash_leaves: dict[str, str],
    readable_leaves: dict[str, Any],
    assumptions: list[AssumptionRecord] | None,
) -> None:
    """Recursively walk a value into flat ``path -> leaf`` dictionaries.

    Handles the exact shapes that appear in this package's inputs and result
    objects: :class:`~.contracts.Assumption`, ``None``, ``bool``, ``int``/
    ``float`` (incl. numpy scalars), ``str``, any ``Mapping`` (incl.
    ``MappingProxyType``), ``numpy.ndarray``/``pandas.Series``, ``list``/
    ``tuple``, and any frozen dataclass (every result type in :mod:`dcf`,
    :mod:`comps` and :mod:`threestatement` is one). Anything else is refused
    rather than silently stringified, because a silent ``str(obj)`` fallback
    could make two different objects hash identically.

    Args:
        value: The (sub)value to flatten.
        path: The dotted/bracket path ``value`` was reached at. Callers
            should start with ``"$"``.
        hash_leaves: Mutated in place: ``path -> canonical string`` for
            hashing.
        readable_leaves: Mutated in place: ``path -> plain Python scalar``
            for display/diffing.
        assumptions: If not ``None``, mutated in place: every
            :class:`~.contracts.Assumption` encountered is appended as an
            :class:`AssumptionRecord`. Pass ``None`` when flattening raw
            inputs (assumptions are recorded from the RESULT only -- see the
            module docstring).

    Raises:
        TypeError: If ``value`` is of a type this function does not know how
            to normalize.
    """
    if isinstance(value, np.generic):
        value = value.item()

    if isinstance(value, Assumption):
        if assumptions is not None:
            assumptions.append(AssumptionRecord(path=path, assumption=value))
        _flatten(value.name, f"{path}.name", hash_leaves, readable_leaves, assumptions)
        _flatten(value.value, f"{path}.value", hash_leaves, readable_leaves, assumptions)
        _flatten(value.basis, f"{path}.basis", hash_leaves, readable_leaves, assumptions)
        _flatten(value.source, f"{path}.source", hash_leaves, readable_leaves, assumptions)
        return

    if value is None:
        hash_leaves[path] = "null"
        readable_leaves[path] = None
        return

    if isinstance(value, bool):
        hash_leaves[path] = "true" if value else "false"
        readable_leaves[path] = value
        return

    if isinstance(value, (int, float)):
        hash_leaves[path] = _normalize_number(value)
        readable_leaves[path] = float(value)
        return

    if isinstance(value, str):
        hash_leaves[path] = json.dumps(value, ensure_ascii=False)
        readable_leaves[path] = value
        return

    if isinstance(value, Mapping):
        for key in sorted(value.keys(), key=str):
            _flatten(value[key], f"{path}.{key}", hash_leaves, readable_leaves, assumptions)
        return

    if isinstance(value, np.ndarray):
        value = value.tolist()
    elif isinstance(value, pd.Series):
        value = value.tolist()

    if isinstance(value, (list, tuple)):
        for index, item in enumerate(value):
            _flatten(item, f"{path}[{index}]", hash_leaves, readable_leaves, assumptions)
        return

    if dataclasses.is_dataclass(value) and not isinstance(value, type):
        for field in dataclasses.fields(value):
            _flatten(
                getattr(value, field.name),
                f"{path}.{field.name}",
                hash_leaves,
                readable_leaves,
                assumptions,
            )
        return

    raise TypeError(
        f"artifact: cannot normalize a value of type {type(value).__name__!r} at "
        f"path {path!r}; extend _flatten or pass a plain mapping / sequence / "
        "dataclass / Assumption / scalar"
    )


def _hash_leaves(hash_leaves: Mapping[str, str]) -> str:
    """Hash a completed ``path -> canonical string`` leaf set.

    Args:
        hash_leaves: Output of one or more :func:`_flatten` calls.

    Returns:
        A 64-character sha256 hex digest.
    """
    canonical = json.dumps(
        {"schema": _HASH_SCHEMA, "leaves": dict(hash_leaves)},
        sort_keys=True,
        ensure_ascii=False,
        separators=(",", ":"),
    )
    return hashlib.sha256(canonical.encode("utf-8")).hexdigest()


def compute_input_hash(payload: Any) -> str:
    """Hash an arbitrary payload the same way ``build_*_artifact`` hashes inputs.

    Public mainly so the normalization scheme (:func:`_flatten` /
    :func:`_normalize_number`) can be tested directly, independent of any one
    model's builder.

    Args:
        payload: Any value :func:`_flatten` can walk -- typically a mapping
            combining call-level config with the model's raw inputs.

    Returns:
        A 64-character sha256 hex digest. Identical for two payloads that are
        the same after normalization (dict key order, ``int``/``float``,
        ``-0.0``/``0.0`` do not matter); different if any leaf's value
        genuinely differs.
    """
    hash_leaves: dict[str, str] = {}
    readable_leaves: dict[str, Any] = {}
    _flatten(payload, "$", hash_leaves, readable_leaves, None)
    return _hash_leaves(hash_leaves)


def _module_version(module: ModuleType) -> str:
    """Build a content-addressed version tag for the module that produced a result.

    Hashes ``module``'s own source together with :mod:`.contracts`'s, since
    every model in this package depends on ``contracts.py``'s input
    discipline -- a change there changes what produced an artifact even if
    ``module``'s own text is untouched. Content-addressed rather than a
    hand-maintained string because no ``__version__`` exists anywhere in
    :mod:`src.quantlib` to defer to, and a hand-maintained tag goes stale
    exactly when it matters most (a formula changes but the tag does not).

    Args:
        module: The model module (:mod:`dcf`, :mod:`comps` or
            :mod:`threestatement`).

    Returns:
        ``"<module.__name__>@sha256:<16 hex chars>"``.
    """
    own_source = Path(module.__file__).read_bytes()
    shared_source = Path(_contracts_module.__file__).read_bytes()
    digest = hashlib.sha256(own_source + b"\x00" + shared_source).hexdigest()[:16]
    return f"{module.__name__}@sha256:{digest}"


# ---------------------------------------------------------------------------
# Gap/exclusion extraction -- the "ugly parts" every export must carry
# ---------------------------------------------------------------------------


def _dcf_gaps(result: DCFResult) -> tuple[str, ...]:
    """Collect DCF-specific notes: GDP-ceiling flag and the non-chosen terminal method.

    Args:
        result: A completed :class:`DCFResult`.

    Returns:
        Human-readable notes, possibly empty.
    """
    notes: list[str] = []
    tv = result.terminal_value
    if tv.growth_exceeds_gdp_ceiling:
        notes.append(
            f"terminal_growth={tv.terminal_growth.value!r} exceeds the "
            f"gdp_growth_ceiling={tv.gdp_growth_ceiling!r} (a flag, not a "
            "rejection) -- review whether a perpetuity growth rate this high "
            "is defensible for this business."
        )
    other_method = (
        "exit_multiple" if result.terminal_value_method == "perpetuity_growth" else "perpetuity_growth"
    )
    notes.append(
        f"terminal_value_method={result.terminal_value_method!r} was discounted "
        f"into enterprise_value; the other method ({other_method!r}) was still "
        f"computed (perpetuity_terminal_value={tv.perpetuity_terminal_value!r}, "
        f"exit_terminal_value={tv.exit_terminal_value!r}, "
        f"implied_ev_ebitda_multiple={tv.implied_ev_ebitda_multiple!r}, "
        f"implied_perpetuity_growth={tv.implied_perpetuity_growth!r}) and is not "
        "silently dropped."
    )
    return tuple(notes)


def _comps_gaps(result: CompsResult, target: TargetCompany) -> tuple[str, ...]:
    """Collect comps-specific notes: excluded multiples and omitted EV-bridge components.

    Args:
        result: A completed :class:`CompsResult`.
        target: The target company the run was built for -- needed because
            the target's own EV-bridge omissions (used inside
            ``equity_value_from_enterprise_value`` while walking each
            EV-multiple's implied value back to equity) are not retained on
            ``CompsResult`` itself.

    Returns:
        Human-readable notes, possibly empty.
    """
    notes: list[str] = []
    for peer in result.peer_multiples:
        if peer.ev_bridge.omitted_components:
            notes.append(
                f"peer {peer.name!r}: EV bridge omitted "
                f"{', '.join(peer.ev_bridge.omitted_components)} (not supplied -- "
                "treated as unknown, not zero)."
            )
        for excl in peer.exclusions:
            notes.append(
                f"peer {peer.name!r}: {excl.multiple_name} excluded -- denominator "
                f"{excl.denominator_value!r} is not positive."
            )
    target_omitted = tuple(
        name
        for name, value in (
            ("minority_interest", target.minority_interest),
            ("preferred_stock", target.preferred_stock),
            ("investments_in_associates", target.investments_in_associates),
        )
        if value is None
    )
    if target_omitted:
        notes.append(
            f"target {target.name!r}: EV bridge omitted {', '.join(target_omitted)} "
            "when walking an implied enterprise value back to implied equity value."
        )
    return tuple(notes)


def _three_statement_gaps(result: ThreeStatementProjection) -> tuple[str, ...]:
    """Collect three-statement notes: non-convergence and revolver plug activity.

    A successfully-returned :class:`ThreeStatementProjection` cannot actually
    contain a non-converged period -- :func:`~.threestatement.project_three_statement`
    raises :class:`~.threestatement.ConvergenceError` before one is
    constructed. The check below is kept anyway (mirroring
    :mod:`threestatement`'s own defensive re-checks) so that if this invariant
    is ever loosened upstream, an artifact built from the result still
    surfaces it rather than silently reporting ``converged=True``.

    Args:
        result: A completed :class:`ThreeStatementProjection`.

    Returns:
        Human-readable notes, possibly empty.
    """
    notes: list[str] = []
    for period in result.periods:
        if not period.converged:
            notes.append(
                f"period {period.period}: interest/revolver circularity DID NOT "
                f"CONVERGE after {period.iterations} iteration(s) -- treat this "
                "period's financing line items as unresolved."
            )
        cfs = period.cash_flow_statement
        if cfs.revolver_draw > 0.0:
            notes.append(
                f"period {period.period}: revolver drew {cfs.revolver_draw!r} to "
                "hold cash at minimum_cash -- financing is not organic."
            )
        if cfs.revolver_repayment > 0.0:
            notes.append(
                f"period {period.period}: revolver repaid {cfs.revolver_repayment!r} "
                "from surplus cash."
            )
    return tuple(notes)


# ---------------------------------------------------------------------------
# Builders
# ---------------------------------------------------------------------------


def build_dcf_artifact(
    *,
    inputs: Mapping[str, Any],
    result: DCFResult,
    generated_at: datetime,
    capital_structure_basis: str,
    discounting_convention: str,
    terminal_value_method: str,
    gdp_growth_ceiling: float = DEFAULT_LONG_RUN_GDP_GROWTH_CEILING,
) -> ModelArtifact:
    """Wrap a completed :func:`~.dcf.run_dcf` call into a versioned artifact.

    ``capital_structure_basis``, ``discounting_convention``,
    ``terminal_value_method`` and ``gdp_growth_ceiling`` are hashed alongside
    ``inputs`` because each one changes ``result`` just as much as a number in
    ``inputs`` would -- an artifact whose hash ignored them could report
    "inputs unchanged" for two runs that priced the same cash flows two
    different ways.

    Assumptions are collected from ``result`` (via ``result.terminal_value``),
    not from ``inputs`` directly -- :func:`~.dcf.run_dcf` passes the same
    :class:`~.contracts.Assumption` objects through unchanged, so walking the
    result finds them without walking the raw inputs a second time.

    Args:
        inputs: The exact mapping passed as ``run_dcf(inputs, ...)``.
        result: The :class:`DCFResult` that call produced.
        generated_at: Timezone-aware timestamp of the run, supplied by the
            caller. See :func:`_require_generated_at`.
        capital_structure_basis: The ``capital_structure_basis`` the call used.
        discounting_convention: The ``discounting_convention`` the call used.
        terminal_value_method: The ``terminal_value_method`` the call used.
        gdp_growth_ceiling: The ``gdp_growth_ceiling`` the call used.

    Returns:
        A ``model_name="dcf"`` :class:`ModelArtifact`.

    Raises:
        TypeError: If ``generated_at`` is not a ``datetime``, or ``result``
            is not a :class:`DCFResult`.
        ValueError: If ``generated_at`` is timezone-naive.
    """
    generated_at = _require_generated_at(generated_at, "build_dcf_artifact")
    if not isinstance(result, DCFResult):
        raise TypeError(
            f"build_dcf_artifact: result must be a DCFResult, got {type(result).__name__}"
        )

    canonical_inputs = {
        "config": {
            "capital_structure_basis": capital_structure_basis,
            "discounting_convention": discounting_convention,
            "terminal_value_method": terminal_value_method,
            "gdp_growth_ceiling": gdp_growth_ceiling,
        },
        "model_inputs": dict(inputs),
    }
    hash_leaves: dict[str, str] = {}
    readable_inputs: dict[str, Any] = {}
    _flatten(canonical_inputs, "$", hash_leaves, readable_inputs, None)
    input_hash = _hash_leaves(hash_leaves)

    readable_outputs: dict[str, Any] = {}
    assumptions: list[AssumptionRecord] = []
    _flatten(result, "$", {}, readable_outputs, assumptions)

    return ModelArtifact(
        model_name="dcf",
        model_version=_module_version(_dcf_module),
        schema_version=_SCHEMA_VERSION,
        generated_at=generated_at,
        input_hash=input_hash,
        inputs=MappingProxyType(readable_inputs),
        outputs=MappingProxyType(readable_outputs),
        assumptions=tuple(assumptions),
        result=result,
        excluded=_dcf_gaps(result),
        warnings=(),
    )


def build_comps_artifact(
    *,
    target: TargetCompany,
    peers: Sequence[PeerCompany],
    calendarisation_policy: str,
    result: CompsResult,
    generated_at: datetime,
) -> ModelArtifact:
    """Wrap a completed :func:`~.comps.run_comps` call into a versioned artifact.

    :mod:`comps` takes no :class:`~.contracts.Assumption` objects anywhere in
    its required inputs -- every field on :class:`TargetCompany` and
    :class:`PeerCompany` is a bare number or a
    :class:`~.comps.FlowMetricPeriods`. The resulting artifact's
    ``assumptions`` is therefore always ``()``; that is a fact about the
    model, not a gap in this builder, and every export says so explicitly
    rather than leaving the section blank.

    Args:
        target: The :class:`TargetCompany` passed to ``run_comps``.
        peers: The peer sequence passed to ``run_comps``.
        calendarisation_policy: The ``calendarisation_policy`` the call used.
        result: The :class:`CompsResult` that call produced.
        generated_at: Timezone-aware timestamp of the run, supplied by the
            caller.

    Returns:
        A ``model_name="comps"`` :class:`ModelArtifact`.

    Raises:
        TypeError: If ``generated_at`` is not a ``datetime``, or ``result``
            is not a :class:`CompsResult`.
        ValueError: If ``generated_at`` is timezone-naive.
    """
    generated_at = _require_generated_at(generated_at, "build_comps_artifact")
    if not isinstance(result, CompsResult):
        raise TypeError(
            f"build_comps_artifact: result must be a CompsResult, got {type(result).__name__}"
        )

    canonical_inputs = {
        "config": {"calendarisation_policy": calendarisation_policy},
        "target": target,
        "peers": list(peers),
    }
    hash_leaves: dict[str, str] = {}
    readable_inputs: dict[str, Any] = {}
    _flatten(canonical_inputs, "$", hash_leaves, readable_inputs, None)
    input_hash = _hash_leaves(hash_leaves)

    readable_outputs: dict[str, Any] = {}
    assumptions: list[AssumptionRecord] = []
    _flatten(result, "$", {}, readable_outputs, assumptions)

    return ModelArtifact(
        model_name="comps",
        model_version=_module_version(_comps_module),
        schema_version=_SCHEMA_VERSION,
        generated_at=generated_at,
        input_hash=input_hash,
        inputs=MappingProxyType(readable_inputs),
        outputs=MappingProxyType(readable_outputs),
        assumptions=tuple(assumptions),
        result=result,
        excluded=_comps_gaps(result, target),
        warnings=result.warnings,
    )


def build_three_statement_artifact(
    *,
    opening: Mapping[str, float],
    drivers: Mapping[str, Sequence[float]],
    result: ThreeStatementProjection,
    generated_at: datetime,
    circularity_tolerance: float = CIRCULARITY_TOLERANCE_REL,
    max_circularity_iterations: int = MAX_CIRCULARITY_ITERATIONS,
    balance_tolerance: float = BALANCE_TOLERANCE_REL,
) -> ModelArtifact:
    """Wrap a completed :func:`~.threestatement.project_three_statement` call into an artifact.

    The three tolerance/iteration-budget parameters default to the same
    constants :func:`~.threestatement.project_three_statement` itself
    defaults to -- these are numerical-solver knobs, not valuation opinions,
    so mirroring the model's own defaults keeps the hash consistent for a
    caller who did not override them, without this module inventing anything
    :mod:`.contracts` would refuse (a terminal growth rate or similar
    price-determining assumption is never defaulted anywhere in this file).

    Like :mod:`comps`, :mod:`threestatement` takes no
    :class:`~.contracts.Assumption` objects -- every driver is a bare number.
    ``assumptions`` on the returned artifact is therefore always ``()``.

    Args:
        opening: The exact mapping passed as ``project_three_statement(opening, ...)``.
        drivers: The exact mapping passed as the ``drivers`` argument.
        result: The :class:`ThreeStatementProjection` that call produced.
        generated_at: Timezone-aware timestamp of the run, supplied by the
            caller.
        circularity_tolerance: The ``circularity_tolerance`` the call used.
        max_circularity_iterations: The ``max_circularity_iterations`` the
            call used.
        balance_tolerance: The ``balance_tolerance`` the call used.

    Returns:
        A ``model_name="three_statement"`` :class:`ModelArtifact`.

    Raises:
        TypeError: If ``generated_at`` is not a ``datetime``, or ``result``
            is not a :class:`ThreeStatementProjection`.
        ValueError: If ``generated_at`` is timezone-naive.
    """
    generated_at = _require_generated_at(generated_at, "build_three_statement_artifact")
    if not isinstance(result, ThreeStatementProjection):
        raise TypeError(
            "build_three_statement_artifact: result must be a "
            f"ThreeStatementProjection, got {type(result).__name__}"
        )

    canonical_inputs = {
        "config": {
            "circularity_tolerance": circularity_tolerance,
            "max_circularity_iterations": max_circularity_iterations,
            "balance_tolerance": balance_tolerance,
        },
        "opening": dict(opening),
        "drivers": {key: list(value) for key, value in drivers.items()},
    }
    hash_leaves: dict[str, str] = {}
    readable_inputs: dict[str, Any] = {}
    _flatten(canonical_inputs, "$", hash_leaves, readable_inputs, None)
    input_hash = _hash_leaves(hash_leaves)

    readable_outputs: dict[str, Any] = {}
    assumptions: list[AssumptionRecord] = []
    _flatten(result, "$", {}, readable_outputs, assumptions)

    return ModelArtifact(
        model_name="three_statement",
        model_version=_module_version(_threestatement_module),
        schema_version=_SCHEMA_VERSION,
        generated_at=generated_at,
        input_hash=input_hash,
        inputs=MappingProxyType(readable_inputs),
        outputs=MappingProxyType(readable_outputs),
        assumptions=tuple(assumptions),
        result=result,
        excluded=_three_statement_gaps(result),
        warnings=(),
    )


# ---------------------------------------------------------------------------
# Diff
# ---------------------------------------------------------------------------


def _diff_leaves(
    baseline: Mapping[str, Any], updated: Mapping[str, Any]
) -> tuple[tuple[FieldChange, ...], tuple[str, ...], tuple[str, ...]]:
    """Diff two flat ``path -> value`` snapshots.

    Args:
        baseline: The "before" flat snapshot.
        updated: The "after" flat snapshot.

    Returns:
        ``(changed, added_paths, removed_paths)`` -- ``changed`` covers only
        paths present on both sides whose value differs; ``added_paths`` /
        ``removed_paths`` are sorted tuples of paths present on one side only.
    """
    baseline_keys = set(baseline.keys())
    updated_keys = set(updated.keys())
    changed: list[FieldChange] = []
    for path in sorted(baseline_keys & updated_keys):
        old, new = baseline[path], updated[path]
        if old == new:
            continue
        delta = None
        if (
            isinstance(old, (int, float))
            and not isinstance(old, bool)
            and isinstance(new, (int, float))
            and not isinstance(new, bool)
        ):
            delta = float(new) - float(old)
        changed.append(FieldChange(path=path, old=old, new=new, delta=delta))
    added = tuple(sorted(updated_keys - baseline_keys))
    removed = tuple(sorted(baseline_keys - updated_keys))
    return tuple(changed), added, removed


def diff_artifacts(baseline: ModelArtifact, updated: ModelArtifact) -> ArtifactDiff:
    """Diff two artifacts of the same model: which fields changed, by how much.

    Args:
        baseline: The "before" artifact.
        updated: The "after" artifact. Must share ``baseline.model_name``.

    Returns:
        An :class:`ArtifactDiff` naming every changed/added/removed input and
        output path.

    Raises:
        ValueError: If ``baseline.model_name != updated.model_name`` -- the
            two artifacts would have no comparable field paths.
    """
    if baseline.model_name != updated.model_name:
        raise ValueError(
            f"diff_artifacts: cannot diff a {baseline.model_name!r} artifact "
            f"against a {updated.model_name!r} artifact -- they have no "
            "comparable field paths"
        )
    changed_inputs, added_inputs, removed_inputs = _diff_leaves(baseline.inputs, updated.inputs)
    changed_outputs, added_outputs, removed_outputs = _diff_leaves(baseline.outputs, updated.outputs)
    return ArtifactDiff(
        model_name=baseline.model_name,
        baseline_generated_at=baseline.generated_at,
        updated_generated_at=updated.generated_at,
        baseline_input_hash=baseline.input_hash,
        updated_input_hash=updated.input_hash,
        input_hash_changed=baseline.input_hash != updated.input_hash,
        baseline_model_version=baseline.model_version,
        updated_model_version=updated.model_version,
        model_version_changed=baseline.model_version != updated.model_version,
        changed_inputs=changed_inputs,
        added_input_paths=added_inputs,
        removed_input_paths=removed_inputs,
        changed_outputs=changed_outputs,
        added_output_paths=added_outputs,
        removed_output_paths=removed_outputs,
    )


# ---------------------------------------------------------------------------
# xlsx export -- shared helpers
# ---------------------------------------------------------------------------


def _bold_row(ws: Worksheet, row: int = 1) -> None:
    """Bold every cell in one worksheet row (typically a header row)."""
    for cell in ws[row]:
        cell.font = Font(bold=True)


def _write_assumptions_and_gaps_sheet(wb: Workbook, artifact: ModelArtifact) -> None:
    """Write the "Assumptions & Data Quality" sheet shared by every xlsx export.

    Layout (rows, 1-indexed):
        1: ``["Provenance"]``
        2-6: ``model_name`` / ``model_version`` / ``schema_version`` /
             ``generated_at`` / ``input_hash``, each as ``[field, value]``
        7: blank
        8: ``["Assumptions", "name", "value", "basis", "source"]`` header
        9..: one row per :class:`AssumptionRecord`, or a single
             ``"(none) -- ..."`` row if there are none
        blank, ``["Warnings"]``, one row per warning (or ``"(none)"``)
        blank, ``["Excluded / Omitted / Data-quality notes"]``, one row per
             note in ``artifact.excluded`` (or ``"(none)"``)

    Args:
        wb: The workbook to add the sheet to.
        artifact: The artifact whose provenance/assumptions/gaps to write.
    """
    ws = wb.create_sheet("Assumptions & Data Quality")
    ws.append(["Provenance"])
    ws.append(["model_name", artifact.model_name])
    ws.append(["model_version", artifact.model_version])
    ws.append(["schema_version", artifact.schema_version])
    ws.append(["generated_at", artifact.generated_at.isoformat()])
    ws.append(["input_hash", artifact.input_hash])
    ws.append([])
    ws.append(["Assumptions", "name", "value", "basis", "source"])
    if artifact.assumptions:
        for record in artifact.assumptions:
            a = record.assumption
            ws.append([record.path, a.name, a.value, a.basis, a.source or ""])
    else:
        ws.append(["(none) -- this model requires no caller-supplied Assumption objects"])
    ws.append([])
    ws.append(["Warnings"])
    for warning in artifact.warnings:
        ws.append([warning])
    if not artifact.warnings:
        ws.append(["(none)"])
    ws.append([])
    ws.append(["Excluded / Omitted / Data-quality notes"])
    for note in artifact.excluded:
        ws.append([note])
    if not artifact.excluded:
        ws.append(["(none)"])

    _bold_row(ws, 1)
    _bold_row(ws, 8)


# ---------------------------------------------------------------------------
# xlsx export -- DCF
# ---------------------------------------------------------------------------

_FCFF_HEADERS: tuple[str, ...] = (
    "Year", "EBIT", "Tax Rate", "NOPAT", "D&A", "CapEx", "Delta NWC",
    "FCFF", "Discount Factor", "Present Value",
)

_WACC_FIELDS: tuple[str, ...] = (
    "risk_free_rate", "beta", "equity_risk_premium", "size_premium",
    "country_risk_premium", "cost_of_equity", "pretax_cost_of_debt", "tax_rate",
    "cost_of_debt_after_tax", "capital_structure_basis", "equity_weight",
    "debt_weight", "wacc",
)


def _write_fcff_bridge_sheet(wb: Workbook, result: DCFResult) -> None:
    """Write the "FCFF Bridge" sheet: one row per year plus a PV total row."""
    ws = wb.create_sheet("FCFF Bridge")
    ws.append(list(_FCFF_HEADERS))
    for year, discounted in zip(result.fcff_bridge, result.discounted_fcff, strict=True):
        ws.append(
            [
                year.year, year.ebit, year.tax_rate, year.nopat,
                year.depreciation_amortization, year.capex, year.delta_nwc, year.fcff,
                discounted.discount_factor, discounted.present_value,
            ]
        )
    ws.append(
        ["Total", None, None, None, None, None, None, None,
         "PV of explicit FCFF", result.pv_of_explicit_fcff]
    )
    _bold_row(ws)


def _write_wacc_sheet(wb: Workbook, result: DCFResult) -> None:
    """Write the "WACC Build" sheet as ``field -> value`` rows."""
    ws = wb.create_sheet("WACC Build")
    ws.append(["Field", "Value"])
    for field in _WACC_FIELDS:
        ws.append([field, getattr(result.wacc_build, field)])
    _bold_row(ws)


def _write_terminal_value_sheet(wb: Workbook, result: DCFResult) -> None:
    """Write the "Terminal Value" sheet: both estimates, cross-checks, the full bridge."""
    ws = wb.create_sheet("Terminal Value")
    ws.append(["Field", "Value"])
    tv = result.terminal_value
    ndb = result.net_debt_bridge
    rows: tuple[tuple[str, Any], ...] = (
        ("wacc_rate", tv.wacc_rate),
        ("terminal_growth.value", tv.terminal_growth.value),
        ("terminal_growth.basis", tv.terminal_growth.basis),
        ("terminal_growth.source", tv.terminal_growth.source or ""),
        ("perpetuity_terminal_value", tv.perpetuity_terminal_value),
        ("terminal_year_ebitda", tv.terminal_year_ebitda),
        ("exit_multiple.value", tv.exit_multiple.value),
        ("exit_multiple.basis", tv.exit_multiple.basis),
        ("exit_multiple.source", tv.exit_multiple.source or ""),
        ("exit_terminal_value", tv.exit_terminal_value),
        ("implied_ev_ebitda_multiple", tv.implied_ev_ebitda_multiple),
        ("implied_perpetuity_growth", tv.implied_perpetuity_growth),
        ("gdp_growth_ceiling", tv.gdp_growth_ceiling),
        ("growth_exceeds_gdp_ceiling", tv.growth_exceeds_gdp_ceiling),
        ("terminal_value_method (chosen)", result.terminal_value_method),
        ("discounting_convention", result.discounting_convention),
        ("terminal_discount_factor", result.terminal_discount_factor),
        ("pv_of_explicit_fcff", result.pv_of_explicit_fcff),
        ("pv_of_terminal_value", result.pv_of_terminal_value),
        ("enterprise_value", result.enterprise_value),
        ("net_debt_bridge.total_debt", ndb.total_debt),
        ("net_debt_bridge.cash_and_equivalents", ndb.cash_and_equivalents),
        ("net_debt_bridge.minority_interest", ndb.minority_interest),
        ("net_debt_bridge.preferred_equity", ndb.preferred_equity),
        ("net_debt_bridge.associate_investments", ndb.associate_investments),
        ("equity_value", result.equity_value),
        ("net_debt_bridge.diluted_shares", ndb.diluted_shares),
        ("value_per_share", result.value_per_share),
    )
    for name, value in rows:
        ws.append([name, value])
    _bold_row(ws)


def _write_sensitivity_sheet(wb: Workbook, grid: pd.DataFrame) -> None:
    """Write the "Sensitivity" sheet from a wacc x growth value-per-share grid."""
    ws = wb.create_sheet("Sensitivity")
    header_label = f"{grid.index.name} \\ {grid.columns.name}"
    ws.append([header_label] + [float(col) for col in grid.columns])
    for wacc_value, row in grid.iterrows():
        cells = [None if pd.isna(v) else float(v) for v in row]
        ws.append([float(wacc_value)] + cells)
    _bold_row(ws)


def export_dcf_workbook(
    artifact: ModelArtifact,
    path: str | Path,
    *,
    sensitivity_grid: pd.DataFrame,
) -> Path:
    """Export a DCF artifact to a 5-sheet workbook.

    Sheets, in order: "FCFF Bridge", "WACC Build", "Terminal Value" (both
    terminal-value estimates cross-checked, and the full net-debt bridge to
    value per share), "Sensitivity" (the caller-supplied grid), "Assumptions
    & Data Quality" (see :func:`_write_assumptions_and_gaps_sheet`).

    Args:
        artifact: A ``model_name="dcf"`` artifact from :func:`build_dcf_artifact`.
        path: Output ``.xlsx`` path. Parent directories are created if absent.
        sensitivity_grid: A wacc x growth value-per-share grid, e.g. from
            :func:`src.quantlib.valuation.dcf.sensitivity_grid`. Required --
            this function does not choose a WACC/growth sweep range on the
            caller's behalf; see the module docstring.

    Returns:
        The written path.

    Raises:
        TypeError: If ``artifact.model_name != "dcf"``.
    """
    if artifact.model_name != "dcf":
        raise TypeError(
            f"export_dcf_workbook: expected a 'dcf' artifact, got {artifact.model_name!r}"
        )
    result: DCFResult = artifact.result

    wb = Workbook()
    wb.remove(wb.active)
    _write_fcff_bridge_sheet(wb, result)
    _write_wacc_sheet(wb, result)
    _write_terminal_value_sheet(wb, result)
    _write_sensitivity_sheet(wb, sensitivity_grid)
    _write_assumptions_and_gaps_sheet(wb, artifact)

    out_path = Path(path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# xlsx export -- comps
# ---------------------------------------------------------------------------

_PEER_HEADERS: tuple[str, ...] = (
    "name", "equity_value", "total_debt", "cash_and_equivalents", "enterprise_value",
    "omitted_ev_bridge_components", "ebitda", "ebit", "revenue", "diluted_eps",
    "ev_ebitda", "ev_ebit", "ev_sales", "pe", "pb",
)


def _write_peer_detail_sheet(wb: Workbook, result: CompsResult) -> None:
    """Write the "Peer Detail" sheet: one row per peer's bridge, metrics and multiples."""
    ws = wb.create_sheet("Peer Detail")
    ws.append(list(_PEER_HEADERS))
    for peer in result.peer_multiples:
        ws.append(
            [
                peer.name,
                peer.ev_bridge.equity_value,
                peer.ev_bridge.total_debt,
                peer.ev_bridge.cash_and_equivalents,
                peer.ev_bridge.enterprise_value,
                ", ".join(peer.ev_bridge.omitted_components) or "(none)",
                peer.calendarised_ebitda.value,
                peer.calendarised_ebit.value,
                peer.calendarised_revenue.value,
                peer.calendarised_diluted_eps.value,
                peer.ev_ebitda if peer.ev_ebitda is not None else "excluded",
                peer.ev_ebit if peer.ev_ebit is not None else "excluded",
                peer.ev_sales if peer.ev_sales is not None else "excluded",
                peer.pe if peer.pe is not None else "excluded",
                peer.pb if peer.pb is not None else "excluded",
            ]
        )
    _bold_row(ws)


def _write_multiple_matrix_sheet(wb: Workbook, result: CompsResult) -> None:
    """Write the "Multiple Matrix" sheet: cross-sectional stats plus excluded-peer detail."""
    ws = wb.create_sheet("Multiple Matrix")
    ws.append(["multiple", "min", "p25", "median", "p75", "max", "n_included", "n_excluded", "warning"])
    for name in MULTIPLE_NAMES:
        dist = result.distributions[name]
        ws.append(
            [
                name, dist.minimum, dist.p25, dist.median, dist.p75, dist.maximum,
                len(dist.included), len(dist.excluded), dist.warning or "",
            ]
        )
    ws.append([])
    ws.append(["Excluded detail: peer", "multiple", "denominator_value"])
    for peer in result.peer_multiples:
        for excl in peer.exclusions:
            ws.append([excl.peer_name, excl.multiple_name, excl.denominator_value])
    _bold_row(ws, 1)


_QUANTILE_LABELS: tuple[str, ...] = ("minimum", "p25", "median", "p75", "maximum")


def _write_implied_valuation_sheet(wb: Workbook, result: CompsResult) -> None:
    """Write the "Implied Valuation" sheet: target metric x peer quantiles, EV and equity."""
    ws = wb.create_sheet("Implied Valuation")
    ws.append(
        [
            "multiple", "is_ev_multiple", "target_metric_value",
            "implied_ev_min", "implied_ev_p25", "implied_ev_median", "implied_ev_p75", "implied_ev_max",
            "implied_equity_min", "implied_equity_p25", "implied_equity_median",
            "implied_equity_p75", "implied_equity_max", "warning",
        ]
    )
    for name in MULTIPLE_NAMES:
        iv = result.implied_valuations[name]
        ev_map = iv.implied_ev_by_quantile
        eq_map = iv.implied_equity_value_by_quantile
        ev_cells = [ev_map.get(label) if ev_map else None for label in _QUANTILE_LABELS]
        eq_cells = [eq_map.get(label) if eq_map else None for label in _QUANTILE_LABELS]
        ws.append([name, iv.is_ev_multiple, iv.target_metric_value, *ev_cells, *eq_cells, iv.warning or ""])
    _bold_row(ws)


def export_comps_workbook(artifact: ModelArtifact, path: str | Path) -> Path:
    """Export a comps artifact to a 4-sheet workbook.

    Sheets, in order: "Peer Detail", "Multiple Matrix" (including the
    excluded-peer detail block), "Implied Valuation", "Assumptions & Data
    Quality" (see :func:`_write_assumptions_and_gaps_sheet` -- for comps this
    sheet's assumptions block always reads "(none)"; see
    :func:`build_comps_artifact`).

    Args:
        artifact: A ``model_name="comps"`` artifact from :func:`build_comps_artifact`.
        path: Output ``.xlsx`` path. Parent directories are created if absent.

    Returns:
        The written path.

    Raises:
        TypeError: If ``artifact.model_name != "comps"``.
    """
    if artifact.model_name != "comps":
        raise TypeError(
            f"export_comps_workbook: expected a 'comps' artifact, got {artifact.model_name!r}"
        )
    result: CompsResult = artifact.result

    wb = Workbook()
    wb.remove(wb.active)
    _write_peer_detail_sheet(wb, result)
    _write_multiple_matrix_sheet(wb, result)
    _write_implied_valuation_sheet(wb, result)
    _write_assumptions_and_gaps_sheet(wb, artifact)

    out_path = Path(path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# xlsx export -- three-statement
# ---------------------------------------------------------------------------

_INCOME_HEADERS: tuple[str, ...] = (
    "period", "revenue", "gross_profit", "opex", "ebitda",
    "depreciation_amortization", "ebit", "interest_expense", "pretax_income",
    "tax_expense", "net_income",
)

_CASHFLOW_HEADERS: tuple[str, ...] = (
    "period", "net_income", "depreciation_amortization",
    "change_in_net_working_capital", "capex", "dividends_paid", "revolver_draw",
    "revolver_repayment", "financing_activities", "net_change_in_cash",
    "cash_flow_from_operations", "free_cash_flow", "converged", "iterations",
)

_BALANCE_HEADERS: tuple[str, ...] = (
    "period", "cash", "net_working_capital", "ppe", "revolver_balance",
    "paid_in_capital", "retained_earnings", "total_assets", "total_liabilities",
    "total_equity", "residual",
)


def _write_income_statement_sheet(wb: Workbook, projection: ThreeStatementProjection) -> None:
    """Write the "Income Statement" sheet: one row per projected period."""
    ws = wb.create_sheet("Income Statement")
    ws.append(list(_INCOME_HEADERS))
    for period in projection.periods:
        s = period.income_statement
        ws.append(
            [
                s.period, s.revenue, s.gross_profit, s.opex, s.ebitda,
                s.depreciation_amortization, s.ebit, s.interest_expense,
                s.pretax_income, s.tax_expense, s.net_income,
            ]
        )
    _bold_row(ws)


def _write_cash_flow_sheet(wb: Workbook, projection: ThreeStatementProjection) -> None:
    """Write the "Cash Flow Statement" sheet, incl. the convergence diagnostics."""
    ws = wb.create_sheet("Cash Flow Statement")
    ws.append(list(_CASHFLOW_HEADERS))
    for period in projection.periods:
        c = period.cash_flow_statement
        ws.append(
            [
                c.period, c.net_income, c.depreciation_amortization,
                c.change_in_net_working_capital, c.capex, c.dividends_paid,
                c.revolver_draw, c.revolver_repayment, c.financing_activities,
                c.net_change_in_cash, c.cash_flow_from_operations, c.free_cash_flow,
                period.converged, period.iterations,
            ]
        )
    _bold_row(ws)


def _write_balance_sheet_sheet(wb: Workbook, projection: ThreeStatementProjection) -> None:
    """Write the "Balance Sheet" sheet, opening plus every period, with the balance residual."""
    ws = wb.create_sheet("Balance Sheet")
    ws.append(list(_BALANCE_HEADERS))
    all_sheets = (projection.opening_balance_sheet, *(p.balance_sheet for p in projection.periods))
    for bs in all_sheets:
        residual = bs.total_assets - bs.total_liabilities - bs.total_equity
        ws.append(
            [
                bs.period, bs.cash, bs.net_working_capital, bs.ppe, bs.revolver_balance,
                bs.paid_in_capital, bs.retained_earnings, bs.total_assets,
                bs.total_liabilities, bs.total_equity, residual,
            ]
        )
    _bold_row(ws)


def export_three_statement_workbook(artifact: ModelArtifact, path: str | Path) -> Path:
    """Export a three-statement artifact to a 4-sheet workbook.

    Sheets, in order: "Income Statement", "Cash Flow Statement" (including
    the per-period ``converged``/``iterations`` circularity diagnostics),
    "Balance Sheet" (including a ``residual`` column, opening balance sheet
    plus every period), "Assumptions & Data Quality" (see
    :func:`_write_assumptions_and_gaps_sheet` -- for three-statement this
    sheet's assumptions block always reads "(none)"; see
    :func:`build_three_statement_artifact`).

    Args:
        artifact: A ``model_name="three_statement"`` artifact from
            :func:`build_three_statement_artifact`.
        path: Output ``.xlsx`` path. Parent directories are created if absent.

    Returns:
        The written path.

    Raises:
        TypeError: If ``artifact.model_name != "three_statement"``.
    """
    if artifact.model_name != "three_statement":
        raise TypeError(
            "export_three_statement_workbook: expected a 'three_statement' "
            f"artifact, got {artifact.model_name!r}"
        )
    projection: ThreeStatementProjection = artifact.result

    wb = Workbook()
    wb.remove(wb.active)
    _write_income_statement_sheet(wb, projection)
    _write_cash_flow_sheet(wb, projection)
    _write_balance_sheet_sheet(wb, projection)
    _write_assumptions_and_gaps_sheet(wb, artifact)

    out_path = Path(path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# pptx export -- one-slide summary, any model
# ---------------------------------------------------------------------------


def _dcf_headline(result: DCFResult) -> list[str]:
    """Return the DCF valuation-conclusion bullet lines."""
    return [
        f"Enterprise value: {result.enterprise_value:,.2f}",
        f"Equity value: {result.equity_value:,.2f}",
        f"Value per share: {result.value_per_share:,.4f}",
        f"WACC: {result.wacc_build.wacc:.4%}",
        f"Terminal value method used: {result.terminal_value_method}",
    ]


def _dcf_sensitivity_lines(grid: pd.DataFrame | None) -> list[str]:
    """Return the DCF sensitivity bullet lines from a caller-supplied grid, if any."""
    if grid is None:
        return ["Sensitivity grid not supplied to this export."]
    values = grid.to_numpy(dtype=float)
    finite = values[np.isfinite(values)]
    if finite.size == 0:
        return [
            "Sensitivity grid supplied but every cell is undefined (growth >= "
            "WACC across the entire swept range)."
        ]
    return [
        f"Value-per-share range across the sensitivity grid: "
        f"{finite.min():,.4f} to {finite.max():,.4f} "
        f"({finite.size} of {values.size} cell(s) defined)"
    ]


def _comps_headline(result: CompsResult) -> list[str]:
    """Return the comps valuation-conclusion bullet lines."""
    lines = [
        f"Target: {result.target_name}",
        f"Calendarisation policy: {result.calendarisation_policy}",
        f"EPS basis: {result.eps_basis}",
    ]
    for name in MULTIPLE_NAMES:
        iv = result.implied_valuations[name]
        if iv.implied_equity_value_by_quantile is None:
            lines.append(f"{name}: no implied valuation (distribution empty)")
        else:
            median = iv.implied_equity_value_by_quantile["median"]
            lines.append(f"{name}: implied equity value at the median peer multiple = {median:,.2f}")
    return lines


def _comps_sensitivity_lines(result: CompsResult) -> list[str]:
    """Return the comps sensitivity bullet lines: each multiple's peer spread."""
    lines: list[str] = []
    for name in MULTIPLE_NAMES:
        dist = result.distributions[name]
        if dist.median is None:
            lines.append(f"{name}: no usable distribution ({dist.warning})")
            continue
        lines.append(
            f"{name}: peer multiple min {dist.minimum:.2f} / median {dist.median:.2f} / "
            f"max {dist.maximum:.2f}"
        )
    return lines


def _three_statement_headline(result: ThreeStatementProjection) -> list[str]:
    """Return the three-statement conclusion bullet lines: the final projected period."""
    last = result.periods[-1]
    return [
        f"Projection horizon: {len(result.periods)} period(s)",
        f"Final period revenue: {last.income_statement.revenue:,.2f}",
        f"Final period net income: {last.income_statement.net_income:,.2f}",
        f"Final period ending cash: {last.balance_sheet.cash:,.2f}",
        f"Final period ending revolver balance: {last.balance_sheet.revolver_balance:,.2f}",
    ]


def _assumption_lines(artifact: ModelArtifact) -> list[str]:
    """Return the "Key Assumptions" bullet lines, or an honest empty-state line."""
    if not artifact.assumptions:
        return ["No caller-supplied Assumption objects were required by this model."]
    return [
        f"{record.assumption.name} = {record.assumption.value!r} ({record.assumption.basis})"
        for record in artifact.assumptions
    ]


def _gap_lines(artifact: ModelArtifact) -> list[str]:
    """Return the "Data Quality Notes" bullet lines, or an honest empty-state line."""
    lines = list(artifact.excluded) + list(artifact.warnings)
    return lines or ["No exclusions, omissions, or warnings were recorded for this run."]


def _write_deck_section(text_frame: TextFrame, heading: str, lines: Sequence[str], *, first: bool = False) -> None:
    """Append one heading paragraph plus one bulleted paragraph per line to a text frame."""
    paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
    paragraph.text = heading
    paragraph.font.bold = True
    paragraph.font.size = Pt(15)
    for line in lines:
        bullet = text_frame.add_paragraph()
        bullet.text = f"• {line}"
        bullet.font.size = Pt(11)
        bullet.level = 1


def export_summary_deck(
    artifact: ModelArtifact,
    path: str | Path,
    *,
    sensitivity_grid: pd.DataFrame | None = None,
) -> Path:
    """Export a one-slide valuation summary: conclusion, assumptions, sensitivity, gaps.

    Works for any ``model_name``. The same provenance footer required on the
    xlsx exports (input hash, model version, generated-at) is written here
    too -- a slide making an unsourced claim is the same failure as an
    unsourced spreadsheet cell.

    Args:
        artifact: Any :class:`ModelArtifact` (``dcf``, ``comps`` or
            ``three_statement``).
        path: Output ``.pptx`` path. Parent directories are created if absent.
        sensitivity_grid: For a ``dcf`` artifact only, the wacc x growth
            value-per-share grid (e.g. from
            :func:`src.quantlib.valuation.dcf.sensitivity_grid`). If omitted
            on a ``dcf`` artifact, the slide says so explicitly rather than
            guessing a range. Ignored for ``comps``/``three_statement``.

    Returns:
        The written path.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.6))
    title_box.text_frame.text = f"{artifact.model_name.upper()} Valuation Summary"

    if artifact.model_name == "dcf":
        headline = _dcf_headline(artifact.result)
        sensitivity = _dcf_sensitivity_lines(sensitivity_grid)
    elif artifact.model_name == "comps":
        headline = _comps_headline(artifact.result)
        sensitivity = _comps_sensitivity_lines(artifact.result)
    else:
        headline = _three_statement_headline(artifact.result)
        sensitivity = ["Sensitivity analysis is not defined for a three-statement projection."]

    body = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9.2), Inches(5.7))
    text_frame = body.text_frame
    text_frame.word_wrap = True
    _write_deck_section(text_frame, "Valuation Conclusion", headline, first=True)
    _write_deck_section(text_frame, "Key Assumptions", _assumption_lines(artifact))
    _write_deck_section(text_frame, "Sensitivity", sensitivity)
    _write_deck_section(text_frame, "Data Quality Notes", _gap_lines(artifact))

    footer = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(9.2), Inches(0.5))
    footer.text_frame.text = (
        f"input_hash={artifact.input_hash} | model_version={artifact.model_version} | "
        f"schema_version={artifact.schema_version} | generated_at={artifact.generated_at.isoformat()}"
    )

    out_path = Path(path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
